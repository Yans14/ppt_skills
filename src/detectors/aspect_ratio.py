import zipfile
from io import BytesIO
from PIL import Image
from pptx import Presentation
from src.models import SlideNode, LayoutError, ErrorType, Severity


def get_native_dimensions(pptx_path: str) -> dict[str, tuple[float, float]]:
    dims = {}
    with zipfile.ZipFile(pptx_path, 'r') as z:
        for name in z.namelist():
            if name.startswith("ppt/media/"):
                try:
                    data = z.read(name)
                    img = Image.open(BytesIO(data))
                    dims[name] = (img.width, img.height)
                except:
                    pass
    return dims


def detect_aspect_ratio_violations(slides: list[SlideNode], pptx_path: str, tolerance: float = 0.05) -> list[LayoutError]:
    errors = []
    native_dims = get_native_dimensions(pptx_path)

    prs = Presentation(pptx_path)
    media_map = {}

    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    blob = shape.image.blob
                    for name, dims in native_dims.items():
                        media_map[str(shape.shape_id)] = dims
                        break
                except:
                    pass

    for slide in slides:
        for img in slide.image_elements:
            if img.id not in media_map:
                continue

            native_w, native_h = media_map[img.id]
            if native_h == 0 or img.rendered_height == 0:
                continue

            native_ratio = native_w / native_h
            rendered_ratio = img.rendered_width / img.rendered_height
            delta = abs(native_ratio - rendered_ratio)

            if delta > tolerance:
                errors.append(LayoutError(
                    type=ErrorType.ASPECT_RATIO,
                    severity=Severity.WARNING,
                    elements=[img.id],
                    message=f"Image distorted (ratio diff: {delta:.2f})"
                ))

    return errors
