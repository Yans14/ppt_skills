from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml
from lxml import etree
from src.models import SlideNode, TextElement, ImageElement, TextStyle, BoundingBox


def emu_to_px(emu: int) -> float:
    return emu / 914400 * 96


def rgb_to_hex(rgb: RGBColor) -> str:
    return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"


def get_theme_colors(presentation: Presentation) -> dict[str, str]:
    """Extract theme colors from the presentation's slide master.
    
    Source: https://stackoverflow.com/a/78501738
    License: CC BY-SA 4.0
    
    Returns a dict mapping theme color names to hex values (without #).
    Example: {'dk1': '000000', 'lt1': 'FFFFFF', 'accent1': '4472C4', ...}
    """
    try:
        theme_part = presentation.slide_master.part.part_related_by(RT.THEME)
        theme = parse_xml(theme_part.blob)
        color_elements = theme.xpath('a:themeElements/a:clrScheme/*')
        result = {}
        for element in color_elements:
            namespaces = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            root = etree.fromstring(etree.tostring(element))
            theme_name = element.tag.replace('{http://schemas.openxmlformats.org/drawingml/2006/main}', '')
            colors = root.xpath('//a:srgbClr/@val', namespaces=namespaces)
            if colors:
                result[theme_name] = colors[0]
            else:
                # Handle system colors (like sysClr for dk1/lt1)
                sys_colors = root.xpath('//a:sysClr/@lastClr', namespaces=namespaces)
                if sys_colors:
                    result[theme_name] = sys_colors[0]
        return result
    except Exception:
        return {}


# Mapping from MSO_THEME_COLOR enum to theme color scheme names
THEME_COLOR_MAP = {
    MSO_THEME_COLOR.DARK_1: 'dk1',
    MSO_THEME_COLOR.LIGHT_1: 'lt1',
    MSO_THEME_COLOR.DARK_2: 'dk2',
    MSO_THEME_COLOR.LIGHT_2: 'lt2',
    MSO_THEME_COLOR.ACCENT_1: 'accent1',
    MSO_THEME_COLOR.ACCENT_2: 'accent2',
    MSO_THEME_COLOR.ACCENT_3: 'accent3',
    MSO_THEME_COLOR.ACCENT_4: 'accent4',
    MSO_THEME_COLOR.ACCENT_5: 'accent5',
    MSO_THEME_COLOR.ACCENT_6: 'accent6',
    MSO_THEME_COLOR.HYPERLINK: 'hlink',
    MSO_THEME_COLOR.FOLLOWED_HYPERLINK: 'folHlink',
}


def safe_get_color_hex(color_format, theme_colors: dict[str, str] | None = None) -> str | None:
    """Safely extract hex color from a ColorFormat object.
    
    Handles RGB colors, scheme colors, and theme colors.
    Returns None if color cannot be extracted.
    """
    if color_format is None:
        return None
    
    try:
        from pptx.enum.dml import MSO_COLOR_TYPE
        
        if color_format.type == MSO_COLOR_TYPE.RGB:
            return rgb_to_hex(color_format.rgb)
        
        elif color_format.type == MSO_COLOR_TYPE.SCHEME:
            # Try to resolve using theme colors if available
            if theme_colors and hasattr(color_format, 'theme_color') and color_format.theme_color:
                theme_key = THEME_COLOR_MAP.get(color_format.theme_color)
                if theme_key and theme_key in theme_colors:
                    return f"#{theme_colors[theme_key].lower()}"
            # Fallback to trying rgb directly
            try:
                return rgb_to_hex(color_format.rgb)
            except AttributeError:
                return None
        
        elif color_format.type == MSO_COLOR_TYPE.THEME:
            if theme_colors and hasattr(color_format, 'theme_color') and color_format.theme_color:
                theme_key = THEME_COLOR_MAP.get(color_format.theme_color)
                if theme_key and theme_key in theme_colors:
                    return f"#{theme_colors[theme_key].lower()}"
            try:
                return rgb_to_hex(color_format.rgb)
            except AttributeError:
                return None
                
    except Exception:
        pass
    
    # Fallback: try direct rgb access
    try:
        return rgb_to_hex(color_format.rgb)
    except (AttributeError, TypeError):
        return None


def parse_presentation(path: str) -> list[SlideNode]:
    prs = Presentation(path)
    slide_width = emu_to_px(prs.slide_width)
    slide_height = emu_to_px(prs.slide_height)
    slides = []
    
    # Extract theme colors once for the entire presentation
    theme_colors = get_theme_colors(prs)

    for idx, slide in enumerate(prs.slides):
        node = SlideNode(index=idx, width=slide_width, height=slide_height)
        
        bg_color = None
        if slide.background.fill.type is not None:
            bg_color = safe_get_color_hex(slide.background.fill.fore_color, theme_colors)
        node.background_color = bg_color

        for shape in slide.shapes:
            shape_id = str(shape.shape_id)
            bbox = BoundingBox(
                x=emu_to_px(shape.left),
                y=emu_to_px(shape.top),
                width=emu_to_px(shape.width),
                height=emu_to_px(shape.height)
            )

            if shape.has_text_frame:
                text = shape.text_frame.text
                style = TextStyle()
                
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            style.font_size = run.font.size.pt
                        if run.font.name:
                            style.font_name = run.font.name
                        if run.font.bold:
                            style.bold = run.font.bold
                        if run.font.color:
                            style.color = safe_get_color_hex(run.font.color, theme_colors)
                        break
                    break

                node.text_elements.append(TextElement(
                    id=shape_id,
                    slide_index=idx,
                    text=text,
                    style=style,
                    bbox=bbox
                ))

            if hasattr(shape, "image"):
                node.image_elements.append(ImageElement(
                    id=shape_id,
                    slide_index=idx,
                    bbox=bbox,
                    rendered_width=bbox.width,
                    rendered_height=bbox.height
                ))

        slides.append(node)

    return slides

